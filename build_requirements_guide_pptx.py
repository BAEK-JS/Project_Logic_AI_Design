# -*- coding: utf-8 -*-
"""요건 문서 작성 가이드 예시 PowerPoint 생성 스크립트."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
def add_title_slide(prs, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def add_bullet_slide(prs, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    tf.word_wrap = True


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "요건 문서 작성 가이드",
        "로직 맵 · AI 다이어그램 입력용 예시 템플릿\n(증권/자본시장 개발)",
    )

    add_bullet_slide(
        prs,
        "이 가이드의 목적",
        [
            "요건을 한 번에 정리해 두면, AI가 전체 흐름(Mermaid)과 구간별 코드 예시를 잘 생성합니다.",
            "‘무엇을’, ‘누가’, ‘어떤 순서로’, ‘예외는 무엇인지’를 문장으로 쓰는 것이 핵심입니다.",
            "본 PPT는 복사하여 회사 양식에 옮기거나, 섹션 제목만 빌려 실제 요건서를 작성할 수 있습니다.",
        ],
    )

    add_bullet_slide(
        prs,
        "좋은 요건의 조건 (AI 친화)",
        [
            "액터(고객, HTS, 주문서버, 거래소 등)와 시스템 경계가 문장으로 드러난다.",
            "처리 순서가 번호 또는 ‘선행/후행’ 관계로 적혀 있다.",
            "입력·출력 데이터(필드 수준이면 더 좋음)가 구간마다 언급된다.",
            "정상 흐름과 예외(거부, 타임아웃, 재시도)가 분리되어 있다.",
            "약어는 최초 등장 시 풀어 쓴다. (예: IOC(Immediate or Cancel))",
        ],
    )

    add_bullet_slide(
        prs,
        "권장 목차 템플릿",
        [
            "1. 개요 · 배경 · 용어 정의",
            "2. 범위(포함/제외) 및 가정",
            "3. 이해관계자 및 인터페이스(대내·대외 시스템)",
            "4. 업무 흐름(정상 시나리오, 단계별)",
            "5. 데이터·메시지(주요 테이블/API/메시지 유형)",
            "6. 예외·오류·재처리 규칙",
            "7. 비기능(성능, 장애, 감사 로그 등) — 필요 시",
            "8. 미결 사항 · 오픈 이슈",
        ],
    )

    add_bullet_slide(
        prs,
        "필수 체크리스트 (붙여넣기 전)",
        [
            "□ 첫 문단에 ‘이 기능이 해결하는 업무 문제’ 한 줄 요약이 있는가?",
            "□ 정상 플로우가 5~15단계로 나눌 수 있을 만큼 구체적인가?",
            "□ ‘만약 ~이면’ 예외가 3가지 이상 적혀 있는가?",
            "□ 외부 연동(거래소, 예탁, 결제 등)이 있으면 호출 방향이 명확한가?",
            "□ 동시성·중복 방지(멱등 키 등) 요구가 있으면 서술했는가?",
        ],
    )

    add_bullet_slide(
        prs,
        "예시 요건 본문 ① — 개요·범위",
        [
            "[가상 사례] 모바일 앱에서 국내 주식 지정가 주문을 접수하고, 거래소 체결 결과를 고객에게 통보한다.",
            "범위: 주문 접수·검증·거래소 송신·체결·잔고 반영·푸시 알림까지. 해외주식·신용은 제외.",
            "가정: 고객은 로그인 및 위험고지 동의 완료. 거래 시간 내에만 신규 주문 가능.",
            "용어: 체결(execution), 미체결(outstanding), 호가(order book) — 내부 용어집과 동일하게 사용.",
        ],
    )

    add_bullet_slide(
        prs,
        "예시 요건 본문 ② — 정상 흐름 (단계)",
        [
            "1) 고객이 종목·수량·가격·유효기간을 입력하고 주문 확정.",
            "2) 앱은 잔고·한도·거래가능 여부를 조회하고 통과 시 주문서버에 주문 ID를 요청.",
            "3) 주문서버는 주문을 저장(상태: 접수) 후 거래소 게이트웨이에 신규 주문 메시지 전송.",
            "4) 거래소로부터 부분/전량 체결 통보를 수신하면 주문 상태·체결내역을 갱신.",
            "5) 잔고·포지션 반영 배치 또는 실시간 처리 후 고객 앱에 체결 알림.",
            "6) 모든 단계는 주문 ID·체결번호로 추적 가능해야 한다.",
        ],
    )

    add_bullet_slide(
        prs,
        "예시 요건 본문 ③ — 예외·오류",
        [
            "거래소 거부(사유코드 XX): 주문 상태 ‘거부’, 고객에게 사유 표시, 재주문은 신규 건으로.",
            "통신 타임아웃: 주문서버는 재조회 또는 ‘불확실’ 상태로 표시하고, 중복 주문 방지 키로 멱등 처리.",
            "부분체결 후 잔량 취소: 잔량만 취소 요청 가능, 이미 체결분은 롤백 불가.",
            "시세 정지·정지종목: 사전 검증에서 차단, 사용자에게 안내 문구 고정.",
        ],
    )

    add_bullet_slide(
        prs,
        "피하면 좋은 표현",
        [
            "‘적당히’, ‘빠르게’, ‘기존과 비슷하게’ — 기준이 없으면 AI도 추측만 합니다.",
            "요구만 나열하고 순서가 없는 긴 문단 — 단계 번호를 붙이세요.",
            "시스템 이름만 있고 역할 설명이 없는 경우 — 한 줄이라도 역할을 적어 주세요.",
        ],
    )

    add_bullet_slide(
        prs,
        "로직 맵 도구에 넣을 때 팁",
        [
            "위 ‘예시 본문 ①~③’을 하나의 텍스트로 이어 붙여 요건 칸에 넣으면 흐름도 품질이 좋아집니다.",
            "언어(C/Java/Python)를 먼저 고른 뒤 AI 생성 또는 프롬프트 복사를 사용하세요.",
            "생성 후 블록 ID와 다이어그램 노드가 맞는지 한 번만 검토하면 됩니다.",
        ],
    )

    out = Path(__file__).resolve().parent.parent / "requirements_guide_sample_ko.pptx"
    prs.save(out)
    print("Saved:", out)


if __name__ == "__main__":
    main()
